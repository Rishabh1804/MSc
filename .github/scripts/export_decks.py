#!/usr/bin/env python3
"""Export every decks/**/slides.html to a .pptx via headless Chromium.

Each slide is rendered at 1920x1080, screenshot, and added as a full-bleed
image to a python-pptx Presentation. Speaker notes (slides.html sibling
notes.md, if present) are attached as the presentation's notes for slide 1.

Run from the repo root after `pip install playwright python-pptx` and
`playwright install --with-deps chromium`. Outputs build/<deck-path>.pptx.
"""

from __future__ import annotations

import http.server
import pathlib
import socket
import socketserver
import sys
import tempfile
import threading
import time
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

REPO_ROOT = Path(__file__).resolve().parents[2]
DECKS_DIR = REPO_ROOT / "decks"
BUILD_DIR = REPO_ROOT / "build"
SLIDE_WIDTH = 1920
SLIDE_HEIGHT = 1080


def find_free_port() -> int:
    with socket.socket() as s:
        s.bind(("127.0.0.1", 0))
        return s.getsockname()[1]


def serve_repo(port: int) -> threading.Thread:
    handler = http.server.SimpleHTTPRequestHandler
    httpd = socketserver.TCPServer(("127.0.0.1", port), handler)
    httpd.allow_reuse_address = True
    thread = threading.Thread(target=httpd.serve_forever, daemon=True)
    thread.start()
    return thread


def deck_url(port: int, slides_path: Path) -> str:
    rel = slides_path.relative_to(REPO_ROOT).as_posix()
    return f"http://127.0.0.1:{port}/{rel}"


def export_one(slides_path: Path, port: int) -> Path:
    rel = slides_path.parent.relative_to(DECKS_DIR)
    pptx_path = BUILD_DIR / f"{rel.as_posix().replace('/', '__')}.pptx"
    pptx_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH / 96)
    prs.slide_height = Inches(SLIDE_HEIGHT / 96)
    blank = prs.slide_layouts[6]

    with sync_playwright() as p:
        browser = p.chromium.launch()
        ctx = browser.new_context(viewport={"width": SLIDE_WIDTH, "height": SLIDE_HEIGHT})
        page = ctx.new_page()
        page.goto(deck_url(port, slides_path), wait_until="networkidle")
        # Let deck-data-fetch.js populate slides.
        page.wait_for_timeout(1500)

        # Start at the first slide (0, 0) — Reveal.next() walks both
        # horizontal and vertical stacks in document order, so this
        # reaches every slide regardless of whether the deck uses
        # vertical stacks. Reveal.slide(idx) with a single index only
        # sets the horizontal slide, missing verticals.
        page.evaluate("() => window.Reveal && window.Reveal.slide(0, 0)")
        page.wait_for_timeout(400)

        with tempfile.TemporaryDirectory() as td:
            slide_idx = 0
            while True:
                shot = Path(td) / f"slide-{slide_idx:03d}.png"
                page.screenshot(path=str(shot), full_page=False)
                slide = prs.slides.add_slide(blank)
                slide.shapes.add_picture(str(shot), 0, 0, prs.slide_width, prs.slide_height)
                slide_idx += 1

                is_last = page.evaluate(
                    "() => !!(window.Reveal && window.Reveal.isLastSlide && window.Reveal.isLastSlide())"
                )
                if is_last:
                    break
                page.evaluate("() => window.Reveal && window.Reveal.next()")
                page.wait_for_timeout(400)

        browser.close()

    notes_md = slides_path.parent / "notes.md"
    if notes_md.exists() and prs.slides:
        notes_text = notes_md.read_text(encoding="utf-8")
        prs.slides[0].notes_slide.notes_text_frame.text = notes_text

    prs.save(str(pptx_path))
    return pptx_path


def main() -> int:
    decks = sorted(DECKS_DIR.glob("**/slides.html"))
    if not decks:
        print("No decks found under decks/**/slides.html", file=sys.stderr)
        return 0

    BUILD_DIR.mkdir(parents=True, exist_ok=True)
    port = find_free_port()
    cwd_before = pathlib.Path.cwd()
    import os
    os.chdir(REPO_ROOT)
    try:
        serve_repo(port)
        time.sleep(0.3)
        for slides_path in decks:
            out = export_one(slides_path, port)
            print(f"exported {slides_path.relative_to(REPO_ROOT)} -> {out.relative_to(REPO_ROOT)}")
    finally:
        os.chdir(cwd_before)

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
